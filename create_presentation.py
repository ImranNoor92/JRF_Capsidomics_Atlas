#!/usr/bin/env python3
"""
JRF Capsidomics Atlas - PowerPoint Presentation Generator
=========================================================

This script generates a professional PowerPoint presentation
outlining the JRF Capsidomics Atlas project plan.

Requirements:
    pip install python-pptx Pillow

Usage:
    python create_presentation.py

Author: JRF Capsidomics Atlas Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pathlib import Path
import os

# Alias for convenience
RgbColor = RGBColor

# Project paths
PROJECT_ROOT = Path(__file__).parent
FIGURES_DIR = PROJECT_ROOT / "figures"
OUTPUT_PATH = PROJECT_ROOT / "JRF_Capsidomics_Atlas_Presentation.pptx"

# Color scheme (professional blue theme)
COLORS = {
    "primary": RgbColor(0x1a, 0x5f, 0x7a),      # Dark teal
    "secondary": RgbColor(0x2d, 0x9c, 0xdb),    # Bright blue
    "accent": RgbColor(0xe7, 0x4c, 0x3c),       # Red accent
    "text_dark": RgbColor(0x2c, 0x3e, 0x50),    # Dark gray
    "text_light": RgbColor(0xff, 0xff, 0xff),   # White
    "background": RgbColor(0xf8, 0xf9, 0xfa),   # Light gray
    "sjr_blue": RgbColor(0x34, 0x98, 0xdb),     # SJR color
    "djr_red": RgbColor(0xe7, 0x4c, 0x3c),      # DJR color
}


def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "JRF Capsidomics Atlas"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "A Curated Map of Jelly-Roll Fold Proteins\nin Viral Capsids"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER
    
    # Author/Date
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Plan & Methodology"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER


def add_background_slide(prs):
    """Slide 2: Background - What is JRF?"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    add_slide_title(slide, "Background: The Jelly-Roll Fold")
    
    # Content boxes
    # Left box - SJR
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.3), Inches(2.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLORS["sjr_blue"]
    
    sjr_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(4), Inches(0.5))
    tf = sjr_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Single Jelly-Roll (SJR)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_light"]
    
    sjr_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(4), Inches(1.8))
    tf = sjr_content.text_frame
    tf.word_wrap = True
    
    bullets = [
        "• 8-stranded β-barrel",
        "• Tangential orientation",
        "• ssDNA & ssRNA viruses",
        "• T=1, T=3, pseudo-T=3"
    ]
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["text_light"]
    
    # Right box - DJR
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLORS["djr_red"]
    
    djr_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(4), Inches(0.5))
    tf = djr_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Double Jelly-Roll (DJR)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_light"]
    
    djr_content = slide.shapes.add_textbox(Inches(5.4), Inches(2.1), Inches(4), Inches(1.8))
    tf = djr_content.text_frame
    tf.word_wrap = True
    
    bullets = [
        "• Tandem fused β-barrels",
        "• Perpendicular orientation",
        "• dsDNA viruses (PRD1-Adeno)",
        "• T=25 to T>100 (giant viruses)"
    ]
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["text_light"]
    
    # Bottom text
    bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(2))
    tf = bottom_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Insight: The jelly-roll fold is the most widespread capsid architecture in the virosphere, spanning all Baltimore classes and host domains."
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["text_dark"]
    
    p2 = tf.add_paragraph()
    p2.text = "\n• Found in parvoviruses, picornaviruses, adenoviruses, and giant viruses"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLORS["text_dark"]
    
    p3 = tf.add_paragraph()
    p3.text = "• DJR likely arose from SJR gene duplication (ancient evolutionary event)"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLORS["text_dark"]


def add_objectives_slide(prs):
    """Slide 3: Project Objectives"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_slide_title(slide, "Project Objectives")
    
    objectives = [
        ("1", "Master Database", "Create jrf_capsidomics_master.csv with comprehensive annotations for all JRF-containing viral proteins"),
        ("2", "High-Confidence Subset", "Curate a validated subset with structural evidence for reliable analysis"),
        ("3", "McKenna-Style Atlas", "Generate summary tables and figures showing JRF distribution across viral families"),
        ("4", "Aguilar-Style Evolution", "Build structural similarity networks and PFAM co-occurrence analysis"),
        ("5", "Evolutionary Schematic", "Propose SJR→DJR transitions and neofunctionalization pathways"),
    ]
    
    y_pos = 1.5
    for num, title, desc in objectives:
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS["secondary"]
        
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.05), Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER
        
        # Title and description
        obj_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(8), Inches(1))
        tf = obj_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORS["text_dark"]
        
        y_pos += 1.1


def add_methodology_overview_slide(prs):
    """Slide 4: Methodology Overview (6-Phase Pipeline)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_slide_title(slide, "Methodology: 6-Phase Pipeline")
    
    phases = [
        ("Phase 1", "Seed Set", "Literature-grounded gold standard"),
        ("Phase 2", "PFAM Mapping", "Domain identification"),
        ("Phase 3", "Expansion", "Database-wide protein search"),
        ("Phase 4", "Annotation", "McKenna-style capsidomics"),
        ("Phase 5", "Structure", "Aguilar-style evolution"),
        ("Phase 6", "Synthesis", "Figures & summary tables"),
    ]
    
    # Draw pipeline as connected boxes
    box_width = 1.3
    box_height = 0.8
    start_x = 0.5
    y_top = 2.0
    y_bottom = 4.0
    
    for i, (phase, title, desc) in enumerate(phases):
        if i < 3:
            x = start_x + i * 3.0
            y = y_top
        else:
            x = start_x + (5 - i) * 3.0
            y = y_bottom
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        if i < 2:
            box.fill.fore_color.rgb = COLORS["sjr_blue"]
        elif i < 4:
            box.fill.fore_color.rgb = COLORS["secondary"]
        else:
            box.fill.fore_color.rgb = COLORS["djr_red"]
        
        # Phase number
        phase_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(box_width), Inches(0.3))
        tf = phase_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.35), Inches(box_width), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER
        
        # Description below
        desc_box = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y + box_height + 0.1), Inches(box_width + 0.6), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["text_dark"]
        p.alignment = PP_ALIGN.CENTER
    
    # Arrows connecting boxes (simplified - horizontal lines)
    # Top row arrows
    for i in range(2):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(start_x + box_width + i * 3.0 + 0.1), 
            Inches(y_top + box_height/2 - 0.15), Inches(0.8), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLORS["text_dark"]
    
    # Down arrow
    down_arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(start_x + 2 * 3.0 + box_width/2 - 0.15),
        Inches(y_top + box_height + 0.1), Inches(0.3), Inches(0.6)
    )
    down_arrow.fill.solid()
    down_arrow.fill.fore_color.rgb = COLORS["text_dark"]
    
    # Bottom row arrows (reversed)
    for i in range(2):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.LEFT_ARROW, Inches(start_x + box_width + (1-i) * 3.0 + 0.1),
            Inches(y_bottom + box_height/2 - 0.15), Inches(0.8), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLORS["text_dark"]


def add_phase12_slide(prs):
    """Slide 5: Phase 1-2 Details"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_slide_title(slide, "Phase 1-2: Building the Foundation")
    
    # Phase 1 box
    p1_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.5), Inches(4.6), Inches(3.5)
    )
    p1_box.fill.solid()
    p1_box.fill.fore_color.rgb = RgbColor(0xeb, 0xf5, 0xfb)
    
    p1_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(0.5))
    tf = p1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 1: Gold-Standard Seed Set"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    
    p1_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.2), Inches(2.5))
    tf = p1_content.text_frame
    tf.word_wrap = True
    
    content = [
        "✓ 30 confirmed JRF capsid proteins",
        "✓ Spans all genome types:",
        "   • ssDNA: AAV, CPV, PCV2, Geminiviruses",
        "   • ssRNA: Picorna, Noda, Tombus",
        "   • dsDNA: Adenovirus, PRD1, NCLDVs",
        "✓ Includes non-capsid JRF proteins",
        "✓ All have PDB structures"
    ]
    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_dark"]
    
    # Phase 2 box
    p2_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.5), Inches(4.6), Inches(3.5)
    )
    p2_box.fill.solid()
    p2_box.fill.fore_color.rgb = RgbColor(0xfd, 0xed, 0xec)
    
    p2_title = slide.shapes.add_textbox(Inches(5.3), Inches(1.6), Inches(4.2), Inches(0.5))
    tf = p2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 2: PFAM Domain Mapping"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["djr_red"]
    
    p2_content = slide.shapes.add_textbox(Inches(5.3), Inches(2.2), Inches(4.2), Inches(2.5))
    tf = p2_content.text_frame
    tf.word_wrap = True
    
    content = [
        "✓ 19 curated JRF PFAM domains",
        "✓ Classification by fold type:",
        "   • SJR capsid: PF00740, PF00729...",
        "   • DJR capsid: PF00608, PF04451...",
        "   • JRF-derived: PF01107 (30K)",
        "✓ Maps seed → PFAM associations",
        "✓ Enables systematic expansion"
    ]
    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_dark"]
    
    # Bottom output note
    out_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.8))
    tf = out_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Outputs: jrf_seed_set.csv • jrf_pfam_master.csv"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLORS["secondary"]
    p.alignment = PP_ALIGN.CENTER


def add_phase34_slide(prs):
    """Slide 6: Phase 3-4 Details"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_slide_title(slide, "Phase 3-4: Database Expansion & Annotation")
    
    # Phase 3
    p3_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.5))
    tf = p3_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 3: PFAM → All Viral Proteins"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    
    p3_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.5), Inches(2))
    tf = p3_content.text_frame
    tf.word_wrap = True
    content = [
        "• Query InterPro/UniProt for each PFAM",
        "• Filter to virus taxonomy (ID: 10239)",
        "• Collect metadata:",
        "  - Accession, organism, taxonomy",
        "  - Protein length, domain boundaries",
        "• Clean: remove fragments, deduplicate"
    ]
    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_dark"]
    
    # Phase 4
    p4_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(0.5))
    tf = p4_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 4: McKenna-Style Annotation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["djr_red"]
    
    p4_content = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.5), Inches(2))
    tf = p4_content.text_frame
    tf.word_wrap = True
    content = [
        "• Add capsidomics fields:",
        "  - capsid_role (MCP/minor/spike)",
        "  - architecture_class (SJR/DJR)",
        "  - t_number (T=1, T=3, pseudo-T=3...)",
        "  - virion_morphology",
        "• Apply evidence rules → confidence levels"
    ]
    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_dark"]
    
    # Schema table
    table_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.4))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Master Database Schema (key columns)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_dark"]
    
    # Add a simple table representation
    schema_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(1))
    tf = schema_box.text_frame
    p = tf.paragraphs[0]
    p.text = "protein_id | organism | family | capsid_role | architecture | t_number | genome_type | evidence_level"
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.font.color.rgb = COLORS["secondary"]
    
    # Output
    out_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(9), Inches(0.4))
    tf = out_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Output: jrf_capsidomics_master.csv • jrf_high_confidence.csv"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLORS["secondary"]
    p.alignment = PP_ALIGN.CENTER


def add_phase5_slide(prs):
    """Slide 7: Phase 5 - Structural Evolution Analysis"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_slide_title(slide, "Phase 5: Aguilar-Style Structural Evolution")
    
    # Three analysis boxes
    analyses = [
        ("Structural Similarity", 
         "• 16 representative structures\n• TM-align pairwise comparison\n• Similarity matrix & heatmap\n• Hierarchical clustering"),
        ("PFAM Co-occurrence", 
         "• Network: nodes = PFAMs\n• Edges = co-occurrence\n• Identify hybrid architectures\n• Domain combination patterns"),
        ("Evolutionary Inference", 
         "• SJR → DJR duplication\n• PRD1-Adeno-NCLDV lineage\n• Neofunctionalization\n• T-number expansion")
    ]
    
    x_positions = [0.3, 3.4, 6.5]
    colors = [COLORS["sjr_blue"], COLORS["secondary"], COLORS["djr_red"]]
    
    for i, (title, content) in enumerate(analyses):
        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x_positions[i]), Inches(1.5), Inches(3.0), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = colors[i]
        
        title_box = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(1.55), Inches(3.0), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER
        
        # Content box
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x_positions[i]), Inches(2.0), Inches(3.0), Inches(2.2)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = RgbColor(0xf8, 0xf9, 0xfa)
        
        text_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(2.1), Inches(2.8), Inches(2.0))
        tf = text_box.text_frame
        tf.word_wrap = True
        for j, line in enumerate(content.split('\n')):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS["text_dark"]
    
    # Key hypothesis box
    hyp_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(1.3)
    )
    hyp_box.fill.solid()
    hyp_box.fill.fore_color.rgb = RgbColor(0xfc, 0xf3, 0xcf)
    
    hyp_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(8.6), Inches(0.4))
    tf = hyp_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Evolutionary Hypothesis"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_dark"]
    
    hyp_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(8.6), Inches(0.7))
    tf = hyp_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The DJR fold arose from SJR gene duplication, enabling larger capsid sizes (T=25+). This lineage shows vertical inheritance: Bacteria (PRD1) → Archaea (STIV) → Eukarya (Adenovirus, NCLDVs)"
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS["text_dark"]


def add_phase6_slide(prs):
    """Slide 8: Phase 6 - Visualization & Deliverables"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_slide_title(slide, "Phase 6: Visualization & Deliverables")
    
    # Figures section
    fig_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.4))
    tf = fig_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 Generated Figures"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    
    figures = [
        "• Architecture distribution (SJR vs DJR)",
        "• Genome type distribution",
        "• T-number distribution",
        "• Genome × Architecture heatmap",
        "• Top virus families chart",
        "• Structural similarity heatmap",
        "• Hierarchical clustering dendrogram"
    ]
    
    fig_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.5), Inches(2.5))
    tf = fig_box.text_frame
    tf.word_wrap = True
    for i, fig in enumerate(figures):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = fig
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_dark"]
    
    # Tables section
    tbl_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(0.4))
    tf = tbl_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 Summary Tables"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["djr_red"]
    
    tables = [
        "• Family overview (counts, structures)",
        "• Architecture summary",
        "• Genome × Architecture matrix",
        "• PFAM co-occurrence network",
        "• Evolutionary transitions",
        "• Final summary statistics"
    ]
    
    tbl_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.5), Inches(2))
    tf = tbl_box.text_frame
    tf.word_wrap = True
    for i, tbl in enumerate(tables):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = tbl
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_dark"]
    
    # Key deliverable box
    deliv_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(1.2)
    )
    deliv_box.fill.solid()
    deliv_box.fill.fore_color.rgb = COLORS["primary"]
    
    deliv_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(8.6), Inches(0.8))
    tf = deliv_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 Primary Deliverable: jrf_capsidomics_master.csv"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_light"]
    
    p2 = tf.add_paragraph()
    p2.text = "A comprehensive, annotated database of all JRF-containing viral proteins with structural and evolutionary context"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLORS["text_light"]


def add_timeline_slide(prs):
    """Slide 9: Project Timeline"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_slide_title(slide, "Project Timeline")
    
    phases = [
        ("Week 1", "Phase 1-2", "Seed set & PFAM mapping", COLORS["sjr_blue"]),
        ("Week 2", "Phase 3", "Database expansion", COLORS["secondary"]),
        ("Week 3", "Phase 4", "Capsidomics annotation", COLORS["secondary"]),
        ("Week 4", "Phase 5", "Structural analysis", COLORS["djr_red"]),
        ("Week 5", "Phase 6", "Visualization & synthesis", COLORS["djr_red"]),
        ("Week 6", "Review", "Validation & manuscript prep", COLORS["primary"]),
    ]
    
    # Timeline bar
    bar_y = 2.5
    bar_height = 0.8
    
    for i, (week, phase, desc, color) in enumerate(phases):
        x = 0.5 + i * 1.55
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(bar_y), Inches(1.45), Inches(bar_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        
        # Week label
        week_box = slide.shapes.add_textbox(Inches(x), Inches(bar_y - 0.4), Inches(1.45), Inches(0.35))
        tf = week_box.text_frame
        p = tf.paragraphs[0]
        p.text = week
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["text_dark"]
        p.alignment = PP_ALIGN.CENTER
        
        # Phase label
        phase_box = slide.shapes.add_textbox(Inches(x), Inches(bar_y + 0.15), Inches(1.45), Inches(0.3))
        tf = phase_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x - 0.1), Inches(bar_y + bar_height + 0.1), Inches(1.65), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS["text_dark"]
        p.alignment = PP_ALIGN.CENTER
    
    # Milestones
    mile_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.4))
    tf = mile_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Milestones"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    
    milestones = [
        "✓ Week 2: Complete database expansion (10,000+ proteins)",
        "✓ Week 4: Structural similarity matrix and clustering complete",
        "✓ Week 5: All figures and summary tables generated",
        "✓ Week 6: Manuscript-ready atlas and evolutionary schematic"
    ]
    
    mile_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(1.5))
    tf = mile_box.text_frame
    tf.word_wrap = True
    for i, m in enumerate(milestones):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = m
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_dark"]


def add_summary_slide(prs):
    """Slide 10: Summary & Next Steps"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_slide_title(slide, "Summary & Impact")
    
    # Summary points
    summary_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.4))
    tf = summary_title.text_frame
    p = tf.paragraphs[0]
    p.text = "What We're Building"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    
    summaries = [
        "📊 Comprehensive JRF protein database spanning all viral families",
        "🔬 McKenna-style capsidomics annotations (role, T-number, morphology)",
        "🧬 Aguilar-style evolutionary analysis (structure > sequence)",
        "📈 Publication-ready figures and summary tables",
        "🔄 Reproducible Python pipeline for future updates"
    ]
    
    sum_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.8))
    tf = sum_box.text_frame
    tf.word_wrap = True
    for i, s in enumerate(summaries):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["text_dark"]
    
    # Impact section
    impact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(9), Inches(2)
    )
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RgbColor(0xe8, 0xf8, 0xf5)
    
    impact_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(8.6), Inches(0.4))
    tf = impact_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Expected Impact"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    
    impacts = [
        "• Resource for viral capsid engineering (AAV, Adeno vectors)",
        "• Framework for understanding capsid evolution across the virosphere",
        "• Foundation for structure-based classification of new viral isolates",
        "• Template for similar capsidomics projects (other fold families)"
    ]
    
    impact_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(8.6), Inches(1.4))
    tf = impact_text.text_frame
    tf.word_wrap = True
    for i, imp in enumerate(impacts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = imp
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_dark"]


def add_thank_you_slide(prs):
    """Slide 11: Thank You / Questions"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    shape.line.fill.background()
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER
    
    # Questions
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER
    
    # Project link
    link_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = link_box.text_frame
    p = tf.paragraphs[0]
    p.text = "github.com/ImranNoor92/JRF_Capsidomics_Atlas"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER


def add_slide_title(slide, title_text):
    """Add a consistent title to a slide."""
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS["primary"]
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_light"]


def create_presentation():
    """Create the complete presentation."""
    
    print("Creating JRF Capsidomics Atlas Presentation...")
    print("=" * 50)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add slides
    print("Adding slides...")
    
    print("  1. Title slide")
    add_title_slide(prs)
    
    print("  2. Background: What is JRF?")
    add_background_slide(prs)
    
    print("  3. Project Objectives")
    add_objectives_slide(prs)
    
    print("  4. Methodology Overview")
    add_methodology_overview_slide(prs)
    
    print("  5. Phase 1-2: Foundation")
    add_phase12_slide(prs)
    
    print("  6. Phase 3-4: Expansion & Annotation")
    add_phase34_slide(prs)
    
    print("  7. Phase 5: Structural Evolution")
    add_phase5_slide(prs)
    
    print("  8. Phase 6: Deliverables")
    add_phase6_slide(prs)
    
    print("  9. Timeline")
    add_timeline_slide(prs)
    
    print("  10. Summary & Impact")
    add_summary_slide(prs)
    
    print("  11. Thank You")
    add_thank_you_slide(prs)
    
    # Save
    prs.save(str(OUTPUT_PATH))
    
    print("=" * 50)
    print(f"✓ Presentation saved to: {OUTPUT_PATH}")
    print(f"  Total slides: 11")
    
    return OUTPUT_PATH


if __name__ == "__main__":
    create_presentation()
